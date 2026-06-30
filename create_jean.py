import os
import customtkinter as ctk
from customtkinter import filedialog
from docling.document_converter import DocumentConverter
from langchain_community.llms import Ollama
from langchain_core.prompts import PromptTemplate
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches

# 전역 변수로 선택된 파일 경로 저장
selected_files = []

def process_files():
    if not selected_files:
        print("선택된 파일이 없습니다.")
        return
    
    persona = persona_input.get("1.0", "end-1c").strip() or "전문 제안서 작성자"
    llm = Ollama(model="gemma4:e4b")

    for file_path in selected_files:
        try:
            converter = DocumentConverter()
            result = converter.convert(file_path)
            document = result.document
            extracted_text = document.export_to_markdown()

            template = f"당신은 {persona}입니다. 다음 문서를 바탕으로 격식 있고 설득력 있는 제안서 요약본을 작성해 주세요.\n\n문서 내용:\n{{content}}"
            prompt = PromptTemplate.from_template(template)
            chain = prompt | llm
            processed_content = chain.invoke({"content": extracted_text})

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = processed_content
            p.font.size = Pt(14)
            
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            prs.save(f"제안서_{base_name}.pptx")
            print(f"완료: {base_name}.pptx")
        except Exception as e:
            print(f"오류 발생 ({file_path}): {e}")

def select_folder():
    global selected_files
    folder_path = filedialog.askdirectory()
    if folder_path:
        selected_files = [os.path.join(folder_path, f) for f in os.listdir(folder_path) if f.lower().endswith(".pdf")]
        status_label.configure(text=f"{len(selected_files)}개의 파일 선택됨")

def select_files():
    global selected_files
    files = filedialog.askopenfilenames(filetypes=[("PDF files", "*.pdf")])
    if files:
        selected_files = list(files)
        status_label.configure(text=f"{len(selected_files)}개의 파일 선택됨")

# GUI 구성
app = ctk.CTk()
app.title("제안서 생성기")
app.geometry("400x550")

ctk.CTkLabel(app, text="페르소나 설정:").pack(pady=(20, 5))
persona_input = ctk.CTkTextbox(app, height=80, width=300)
persona_input.insert("1.0", "전문 제안서 작성자")
persona_input.pack(pady=10)

# "PDF 선택" 박스
pdf_frame = ctk.CTkFrame(app, fg_color="transparent", border_width=1, border_color="gray")
pdf_frame.pack(pady=20, padx=20, fill="x")
ctk.CTkLabel(pdf_frame, text="PDF 선택").pack(pady=5)
ctk.CTkButton(pdf_frame, text="폴더 단위 선택", command=select_folder).pack(pady=5, padx=10)
ctk.CTkButton(pdf_frame, text="파일 개별 선택", command=select_files).pack(pady=5, padx=10)
status_label = ctk.CTkLabel(pdf_frame, text="파일을 선택하세요", text_color="gray")
status_label.pack(pady=5)

# 실행 버튼
btn_run = ctk.CTkButton(app, text="제안서 생성하기", fg_color="green", command=process_files)
btn_run.pack(pady=20)

app.mainloop()